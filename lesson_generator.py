from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


class Lesson6_Examples_Instead_Of_Images:
    """
    មេរៀនទី ៦៖ ប្តូរពី រូបភាព -> ឧទាហរណ៍ (Example Sentences)
    """
    
    COLORS = {
        'primary': RGBColor(0, 51, 102),      # Navy Blue
        'accent': RGBColor(245, 130, 32),     # Orange
        'text': RGBColor(33, 37, 41),         # Dark Grey
        'white': RGBColor(255, 255, 255),
        'light_blue': RGBColor(235, 245, 255),
        'gray': RGBColor(200, 200, 200),
        'green_excel': RGBColor(33, 115, 70),
        'grid_line': RGBColor(192, 192, 192),
        'trace_color': RGBColor(211, 211, 211)
    }

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def set_font(self, run, size=18, is_title=False, color=None, is_bold=False, font_name=None):
        """កំណត់ font សម្រាប់អត្ថបទ"""
        if font_name:
            run.font.name = font_name
        elif run.font.name and "Microsoft YaHei" in run.font.name:
            pass  # រក្សា Chinese font
        else:
            run.font.name = 'Khmer OS Moul Light' if is_title else 'Khmer OS Battambang'
        
        run.font.size = Pt(size)
        run.font.bold = is_bold
        if color:
            run.font.color.rgb = color

    def set_chinese_font(self, run, size=18, is_bold=True, color=None):
        """កំណត់ font ចិន"""
        run.font.name = 'Microsoft YaHei'
        run.font.size = Pt(size)
        run.font.bold = is_bold
        if color:
            run.font.color.rgb = color

    def add_header(self, slide, title_cn, title_km):
        """បង្កើត header ស្លាយ"""
        # Background header
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), 
            Inches(0), 
            self.prs.slide_width, 
            Inches(1.2)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['primary']
        bg.line.fill.background()  # លុបបន្ទាត់ជុំវិញ
        
        # Title textbox
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title_cn
        for run in p.runs:
            self.set_chinese_font(run, 28, True, self.COLORS['white'])
        
        # Subtitle (Khmer)
        p2 = tb.text_frame.add_paragraph()
        p2.text = title_km
        for run in p2.runs:
            self.set_font(run, 16, is_title=True, color=self.COLORS['white'])

    def draw_tianzi_ge(self, slide, x, y, size, char=""):
        """គូរក្រឡាហាត់សរសេរអក្សរចិន"""
        # ប្រអប់ខាងក្រៅ
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
        box.fill.background()
        box.line.color.rgb = self.COLORS['primary']
        box.line.width = Pt(1.5)

        # បន្ទាត់បញ្ឈរកណ្តាល
        v_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 
            x + size/2, y, 
            x + size/2, y + size
        )
        v_line.line.color.rgb = self.COLORS['grid_line']
        v_line.line.dash_style = 4
        v_line.line.width = Pt(0.5)
        
        # បន្ទាត់ផ្ដេកកណ្តាល
        h_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 
            x, y + size/2, 
            x + size, y + size/2
        )
        h_line.line.color.rgb = self.COLORS['grid_line']
        h_line.line.dash_style = 4
        h_line.line.width = Pt(0.5)

        # ដាក់អក្សរគំរូ (ប្រសិនបើមាន)
        if char:
            tb = slide.shapes.add_textbox(x, y + Inches(0.05), size, size)
            p = tb.text_frame.paragraphs[0]
            p.text = char
            p.alignment = PP_ALIGN.CENTER
            tb.text_frame.vertical_anchor = 1  # កណ្តាលបញ្ឈរ
            for run in p.runs:
                run.font.name = 'KaiTi'
                run.font.size = Pt(42)
                run.font.color.rgb = self.COLORS['trace_color']

    def create_cover(self):
        """បង្កើតស្លាយគម្រប"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            self.prs.slide_width, 
            self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['light_blue']
        bg.line.fill.background()
        
        # ប្រអប់កណ្តាល
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(3), Inches(2), 
            Inches(7.333), Inches(3.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS['white']
        box.line.color.rgb = self.COLORS['primary']
        box.line.width = Pt(3)
        
        # អត្ថបទ
        tb = slide.shapes.add_textbox(Inches(3.2), Inches(2.5), Inches(6.9), Inches(2.5))
        
        # ចំណងជើងចិន
        p = tb.text_frame.paragraphs[0]
        p.text = "第六课：成车异常与 Excel 计数"
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            self.set_chinese_font(run, 32, True, self.COLORS['primary'])
        
        # ចំណងជើងខ្មែរ
        p2 = tb.text_frame.add_paragraph()
        p2.text = "មេរៀនទី ៦៖ បញ្ហាកង់សម្រេច និង រូបមន្ត Excel (COUNTIF)"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
        for run in p2.runs:
            self.set_font(run, 20, is_title=True, color=self.COLORS['text'])
        
        # ឈ្មោះគ្រូបង្រៀន
        p3 = tb.text_frame.add_paragraph()
        p3.text = "培训教师：郑和"
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(30)
        for run in p3.runs:
            self.set_chinese_font(run, 16, True, self.COLORS['accent'])

    def create_vocab_slide(self, title_cn, title_km, vocab_list):
        """បង្កើតស្លាយពាក្យ ជាមួយឧទាហរណ៍"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, title_cn, title_km)
        
        # ចំណងជើងតារាង
        headers = ["中文", "拼音", "ភាសាខ្មែរ", "例句 (ឧទាហរណ៍)"]
        widths = [2.3, 2.3, 2.8, 5.0]  # កែទំហំឱ្យសមរម្យ
        left = Inches(0.4)
        top = Inches(1.4)
        
        current_x = left
        for h, w in zip(headers, widths):
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                current_x, top, 
                Inches(w), Inches(0.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS['primary']
            box.line.fill.background()
            
            tb = slide.shapes.add_textbox(current_x, top, Inches(w), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = h
            p.alignment = PP_ALIGN.CENTER
            tb.text_frame.vertical_anchor = 1  # កណ្តាលបញ្ឈរ
            
            # កំណត់ font ឱ្យត្រឹមត្រូវ
            if "中文" in h or "拼音" in h or "例句" in h:
                for run in p.runs:
                    self.set_chinese_font(run, 12, True, self.COLORS['white'])
            else:
                for run in p.runs:
                    self.set_font(run, 12, is_title=True, color=self.COLORS['white'])
            
            current_x += Inches(w)

        # ជួរឈរពាក្យ
        row_height = Inches(1.6)
        gap = Inches(0.1)
        
        for idx, (cn, py, km, ex_cn, ex_km) in enumerate(vocab_list):
            y = top + Inches(0.6) + (row_height + gap) * idx
            
            # Background row
            bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                left, y, 
                sum([Inches(x) for x in widths]), 
                row_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.COLORS['light_blue'] if idx % 2 == 0 else self.COLORS['white']
            bg.line.color.rgb = self.COLORS['gray']
            bg.line.width = Pt(0.5)
            
            # កំណត់ទីតាំងជួរឈរ
            x_cn = left
            x_py = left + Inches(widths[0])
            x_km = left + Inches(widths[0] + widths[1])
            x_ex = left + Inches(widths[0] + widths[1] + widths[2])

            # ពាក្យចិន
            tb = slide.shapes.add_textbox(x_cn, y + Inches(0.5), Inches(widths[0]), Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            p.text = cn
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                self.set_chinese_font(run, 24, True, self.COLORS['primary'])
            
            # ពិន្យិន
            tb = slide.shapes.add_textbox(x_py, y + Inches(0.6), Inches(widths[1]), Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            p.text = py
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = 'Arial'
                run.font.size = Pt(15)
                run.font.color.rgb = self.COLORS['text']
            
            # ពាក្យខ្មែរ
            tb = slide.shapes.add_textbox(x_km, y + Inches(0.55), Inches(widths[2]), Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            p.text = km
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                self.set_font(run, 17, is_title=False, color=self.COLORS['text'])
            
            # ឧទាហរណ៍ប្រយោគ
            tb_ex = slide.shapes.add_textbox(
                x_ex + Inches(0.15), 
                y + Inches(0.25), 
                Inches(widths[3] - 0.3), 
                Inches(1.2)
            )
            tb_ex.text_frame.word_wrap = True
            
            # ប្រយោគចិន
            p = tb_ex.text_frame.paragraphs[0]
            p.text = ex_cn
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                self.set_chinese_font(run, 13, False, self.COLORS['primary'])
            
            # ប្រយោគខ្មែរ
            p2 = tb_ex.text_frame.add_paragraph()
            p2.text = ex_km
            p2.space_before = Pt(6)
            for run in p2.runs:
                self.set_font(run, 12, False, self.COLORS['text'])

    def create_excel_countif_slide(self):
        """បង្កើតស្លាយ Excel COUNTIF"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "2. Excel 公式：计数 (COUNTIF)", "រូបមន្តរាប់ចំនួនតាមលក្ខខណ្ឌ")
        
        # ប្រអប់ពន្យល់
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.5), Inches(1.6), 
            Inches(5.5), Inches(2.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS['light_blue']
        box.line.color.rgb = self.COLORS['primary']
        box.line.width = Pt(2)
        
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5), Inches(2.4))
        
        # ចំណងជើង
        p = tb.text_frame.paragraphs[0]
        p.text = "🔢 COUNTIF"
        for run in p.runs:
            self.set_font(run, 26, False, self.COLORS['primary'], True, font_name='Arial')
        
        # ពន្យល់ខ្មែរ
        p2 = tb.text_frame.add_paragraph()
        p2.text = "ប្រើសម្រាប់រាប់ចំនួនតាមលក្ខខណ្ឌដែលយើងចង់បាន។"
        p2.space_before = Pt(12)
        for run in p2.runs:
            self.set_font(run, 15, False, self.COLORS['text'])

        # ឧទាហរណ៍
        p3 = tb.text_frame.add_paragraph()
        p3.text = "ឧទាហរណ៍៖ រាប់មើលថាមាន \"NG\" ប៉ុន្មាន?"
        p3.space_before = Pt(15)
        for run in p3.runs:
            self.set_font(run, 14, False, self.COLORS['text'])
        
        # រូបមន្ត
        p4 = tb.text_frame.add_paragraph()
        p4.text = '=COUNTIF(C2:C10, "NG")'
        p4.space_before = Pt(10)
        for run in p4.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = self.COLORS['green_excel']

        # ប្រអប់រូបភាព Excel
        img_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(6.5), Inches(1.6), 
            Inches(6.3), Inches(5.3)
        )
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = self.COLORS['white']
        img_box.line.color.rgb = self.COLORS['gray']
        img_box.line.dash_style = 2
        img_box.line.width = Pt(2)
        
        # អត្ថបទណែនាំ
        tb = slide.shapes.add_textbox(Inches(7), Inches(3.5), Inches(5.3), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "📸 Paste Excel Screenshot Here\n(បញ្ចូលរូបភាពបង្ហាញរូបមន្ត COUNTIF)"
        p.alignment = PP_ALIGN.CENTER
        tb.text_frame.vertical_anchor = 1
        for run in p.runs:
            self.set_font(run, 14, False, self.COLORS['gray'])

    def create_homework(self):
        """បង្កើតស្លាយកិច្ចការផ្ទះ"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "3. 本周作业 (Homework)", "កិច្ចការផ្ទះ")
        
        # ប្រអប់មាតិកា
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(2), Inches(2.5), 
            Inches(9.333), Inches(3.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['light_blue']
        bg.line.color.rgb = self.COLORS['primary']
        bg.line.width = Pt(3)
        
        tb = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(8.333), Inches(2.8))
        
        # ចំណងជើង
        p = tb.text_frame.paragraphs[0]
        p.text = "💻 任务 (Task):"
        for run in p.runs:
            self.set_chinese_font(run, 26, True, self.COLORS['accent'])
        
        # បញ្ជីកិច្ចការ
        p2 = tb.text_frame.add_paragraph()
        p2.text = "1. 抄写生词 12 个（每个字写一行）"
        p2.space_before = Pt(25)
        for run in p2.runs:
            self.set_chinese_font(run, 18, False, self.COLORS['text'])
        
        p3 = tb.text_frame.add_paragraph()
        p3.text = "   សរសេរពាក្យថ្មី ១២ ពាក្យ (គ្រប់តួអក្សរសរសេរម្តង)"
        p3.space_before = Pt(8)
        for run in p3.runs:
            self.set_font(run, 15, False, self.COLORS['text'])
        
        p4 = tb.text_frame.add_paragraph()
        p4.text = "2. 使用 COUNTIF 统计报表中的 NG 数量"
        p4.space_before = Pt(20)
        for run in p4.runs:
            self.set_chinese_font(run, 18, False, self.COLORS['text'])
        
        p5 = tb.text_frame.add_paragraph()
        p5.text = "   ប្រើរូបមន្ត COUNTIF រាប់ចំនួន NG ក្នុងរបាយការណ៍"
        p5.space_before = Pt(8)
        for run in p5.runs:
            self.set_font(run, 15, False, self.COLORS['text'])

    def create_writing_practice_auto(self, lesson_words):
        """បង្កើតតារាងហាត់សរសេរស្វ័យប្រវត្តិ"""
        words_per_page = 14
        chunks = [lesson_words[i:i + words_per_page] 
                  for i in range(0, len(lesson_words), words_per_page)]
        
        for i, chunk in enumerate(chunks):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self.add_header(slide, f"附录 {i+1}：汉字书写练习", "តារាងហាត់សរសេរអក្សរចិន")
            
            start_x = Inches(0.5)
            start_y = Inches(1.6)
            box_size = Inches(0.8)
            gap = Inches(0.08)
            current_y = start_y
            
            for char in chunk:
                # ប្រអប់ទី១ មានតួអក្សរគំរូ
                self.draw_tianzi_ge(slide, start_x, current_y, box_size, char)
                
                # ប្រអប់ទី២-១៤ ទទេសម្រាប់សរសេរ
                for col in range(1, 14):
                    self.draw_tianzi_ge(
                        slide, 
                        start_x + (col * (box_size + Inches(0.02))), 
                        current_y, 
                        box_size, 
                        ""
                    )
                
                current_y += (box_size + gap)

    def generate(self, filename="Lesson_06_Examples.pptx"):
        """បង្កើត PowerPoint ពេញលេញ"""
        
        # ស្លាយគម្រប
        self.create_cover()
        
        # បញ្ជីពាក្យជាមួយឧទាហរណ៍
        vocab1 = [
            ("刹车失灵", "shā chē shī líng", "ហ្វ្រាំងមិនស៊ី", 
             "后轮刹车失灵，很危险。", "ហ្វ្រាំងក្រោយមិនស៊ីទេ គ្រោះថ្នាក់ណាស់។"),
            ("变速不准", "biàn sù bù zhǔn", "ដូរលេខមិនចូល", 
             "这辆车变速不准，需要调试。", "ឡាននេះដូរលេខមិនចូលទេ ត្រូវសារ៉េ។"),
            ("轮胎漏气", "lún tāi lòu qì", "សំបកកង់ធ្លាយ", 
             "前轮漏气了，请更换内胎។", "កង់មុខធ្លាយខ្យល់ហើយ សុំដូរពោះវៀនកង់។")
        ]
        
        vocab2 = [
            ("螺丝松动", "luó sī sōng dòng", "ខ្ចៅធូរ", 
             "脚踏螺丝松动，请锁紧。", "ខ្ចៅជើងធាក់ធូរហើយ សូមរឹតឱ្យតឹង។"),
            ("异响", "yì xiǎng", "សំឡេងរំខាន", 
             "骑行时有异响។", "ពេលជិះមានសំឡេងរំខាន។"),
            ("划痕", "huá hén", "ស្នាមឆ្កូត", 
             "车架上有划痕，是NG品។", "នៅលើតួកង់មានស្នាមឆ្កូត គឺជាផលិតផល NG។")
        ]
        
        vocab3 = [
            ("掉漆", "diào qī", "របកថ្នាំ", 
             "这里掉漆了，需要补漆។", "កន្លែងនេះរបកថ្នាំហើយ ត្រូវការបាញ់ថ្នាំបន្ថែម។"),
            ("生锈", "shēng xiù", "ច្រែះ", 
             "链条生锈了，不能出货។", "ច្រវាក់ឡើងច្រែះហើយ ចេញទំនិញមិនបានទេ។"),
            ("错件", "cuò jiàn", "ដាក់គ្រឿងខុស", 
             "注意不要装错件。", "ប្រយ័ត្ន! កុំដំឡើងគ្រឿងខុស។")
        ]
        
        vocab4 = [
            ("漏装", "lòu zhuāng", "ភ្លេចដាក់គ្រឿង", 
             "你漏装了一个垫片。", "អ្នកភ្លេចដាក់កងមួយ។"),
            ("歪斜", "wāi xié", "វៀច / មិនត្រង់", 
             "车把歪斜，请校正。", "ដៃកង់វៀចហើយ សូមកែតម្រូវ។"),
            ("返工", "fǎn gōng", "ធ្វើឡើងវិញ", 
             "这批货全部需要返工។", "ទំនិញមួយឡូត៍នេះត្រូវធ្វើឡើងវិញទាំងអស់។")
        ]
        
        # បង្កើតស្លាយពាក្យ
        self.create_vocab_slide("1.1 常见异常 (Part 1)", "បញ្ហាទូទៅ ១", vocab1)
        self.create_vocab_slide("1.2 常见异常 (Part 2)", "បញ្ហាទូទៅ ២", vocab2)
        self.create_vocab_slide("1.3 常见异常 (Part 3)", "បញ្ហាទូទៅ ៣", vocab3)
        self.create_vocab_slide("1.4 常见异常 (Part 4)", "បញ្ហាទូទៅ ៤", vocab4)
        
        # ស្លាយ Excel
        self.create_excel_countif_slide()
        
        # ស្លាយកិច្ចការផ្ទះ
        self.create_homework()
        
        # ប្រមូលតួអក្សរទាំងអស់
        all_chars = []
        for v_list in [vocab1, vocab2, vocab3, vocab4]:
            for item in v_list:
                word = item[0]
                for char in word:
                    all_chars.append(char)
        
        # បង្កើតតារាងហាត់សរសេរ
        self.create_writing_practice_auto(all_chars)
        
        # រក្សាទុកឯកសារ
        self.prs.save(filename)
        print(f"✅ បានបង្កើតមេរៀនទី ៦ (១២ ពាក្យ + ឧទាហរណ៍) ជោគជ័យ: {filename}")


if __name__ == "__main__":
    app = Lesson6_Examples_Instead_Of_Images()
    app.generate()